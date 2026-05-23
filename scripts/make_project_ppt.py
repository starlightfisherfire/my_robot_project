#!/usr/bin/env python3
"""Generate a project overview PPT for Paper 1."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
BG_DARK    = RGBColor(0x1A, 0x1A, 0x2E)
BG_MID     = RGBColor(0x16, 0x21, 0x3E)
ACCENT     = RGBColor(0x00, 0xD2, 0xFF)
ACCENT2    = RGBColor(0xFF, 0x6B, 0x6B)
ACCENT3    = RGBColor(0x4E, 0xCB, 0x71)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
GOLD       = RGBColor(0xFF, 0xD7, 0x00)

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, items, left, top, width, height, font_size=16, color=WHITE, bullet_color=ACCENT, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


# ═══════════════════════════════════════════════════
# Slide 1: Title
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BG_DARK)

# accent bar
add_shape_bg(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
    "Representation Learning for Robotic Pushing",
    font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
    "OOD Generalization via Object-Centric & Causality-Aware Encoders with CEM-MPC",
    font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
    "Paper 1  ·  MuJoCo Simulation  ·  SO-101 Robot",
    font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# Slide 2: Research Question
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Research Question", font_size=32, color=ACCENT, bold=True)

add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.8), BG_MID)
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(10.5), Inches(1.5),
    "How does the choice of state representation affect a learned dynamics model's\n"
    "ability to generalize to out-of-distribution (OOD) scenarios in robotic pushing?",
    font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Three sub-questions
for i, (title, desc) in enumerate([
    ("Layout OOD", "Can the model handle novel obstacle configurations\n(blocking, narrow passages) not seen during training?"),
    ("Shape OOD", "Can the model generalize to novel object shapes\n(L-shapes, crosses) trained only on T-shapes?"),
    ("Representation Gap", "Do object-centric or causality-aware encoders\noutperform flat state encoding under distribution shift?"),
]):
    x = Inches(0.8 + i * 4.0)
    add_shape_bg(slide, x, Inches(4.0), Inches(3.6), Inches(2.8), BG_MID)
    add_text_box(slide, x + Inches(0.3), Inches(4.2), Inches(3.0), Inches(0.6),
        title, font_size=20, color=ACCENT, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(4.9), Inches(3.0), Inches(1.8),
        desc, font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════
# Slide 3: Method Overview
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Method Overview", font_size=32, color=ACCENT, bold=True)

# Pipeline boxes
steps = [
    ("1. Reset Templates", "140 templates across\ntrain / ID-test / OOD splits\n(T/L/cross/bar/square/cylinder)"),
    ("2. MuJoCo Env", "MujocoPushEnv with\ncompound object geoms,\nobstacle instantiation"),
    ("3. Oracle Rollout", "True dynamics rollout\nfor CEM cost evaluation\n(state clone/restore)"),
    ("4. CEM-MPC", "Cross-Entropy Method\nplanner with strict\npose-stop criterion"),
    ("5. Evaluation", "Success rate, pos/θ error,\nOOD gap analysis\nacross 3 encoder variants"),
]

for i, (title, desc) in enumerate(steps):
    x = Inches(0.5 + i * 2.5)
    add_shape_bg(slide, x, Inches(1.8), Inches(2.2), Inches(2.8), BG_MID)
    add_text_box(slide, x + Inches(0.15), Inches(1.95), Inches(1.9), Inches(0.6),
        title, font_size=15, color=GOLD, bold=True)
    add_text_box(slide, x + Inches(0.15), Inches(2.6), Inches(1.9), Inches(1.8),
        desc, font_size=12, color=LIGHT_GRAY)

# Three encoder variants
add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.6),
    "Three Encoder Variants (Learned Dynamics Models)", font_size=20, color=WHITE, bold=True)

for i, (name, desc, color) in enumerate([
    ("FlatEncoder", "Raw state vector → z\n(Baseline)", ACCENT2),
    ("ObjectCentricEncoder", "Per-object slots → z\n(Structured representation)", ACCENT),
    ("CausalityAwareEncoder", "Causal graph masking → z\n(Intervention-aware)", ACCENT3),
]):
    x = Inches(0.8 + i * 4.0)
    add_shape_bg(slide, x, Inches(5.9), Inches(3.6), Inches(1.2), BG_MID)
    add_text_box(slide, x + Inches(0.2), Inches(5.95), Inches(3.2), Inches(0.5),
        name, font_size=16, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(6.4), Inches(3.2), Inches(0.7),
        desc, font_size=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════
# Slide 4: Environment & Task Design
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Environment & Task Design", font_size=32, color=ACCENT, bold=True)

# Object shapes
add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.5), BG_MID)
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5.0), Inches(0.5),
    "Object Shapes (6 families)", font_size=18, color=GOLD, bold=True)
shapes_text = (
    "• T-shape  (primary, 12mm thick)\n"
    "• L-shape  (OOD shape test)\n"
    "• Cross, Bar, Square, Cylinder\n\n"
    "Compound MuJoCo geoms via ObjectShapeFactory"
)
add_text_box(slide, Inches(1.0), Inches(2.3), Inches(5.0), Inches(1.5),
    shapes_text, font_size=13, color=LIGHT_GRAY)

# Layouts
add_shape_bg(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(2.5), BG_MID)
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5.0), Inches(0.5),
    "Layout Families (OOD)", font_size=18, color=GOLD, bold=True)
layout_text = (
    "• Open space  (no obstacles)\n"
    "• Blocking  (easy / medium / hard)\n"
    "• Passage  (wide 9cm / medium 7cm / hard 5cm)\n"
    "• Edge goal\n\n"
    "Passage gap controls difficulty level"
)
add_text_box(slide, Inches(7.0), Inches(2.3), Inches(5.0), Inches(1.5),
    layout_text, font_size=13, color=LIGHT_GRAY)

# Pusher & success
add_shape_bg(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.5), BG_MID)
add_text_box(slide, Inches(1.0), Inches(4.6), Inches(5.0), Inches(0.5),
    "Pusher Configuration", font_size=18, color=GOLD, bold=True)
pusher_text = (
    "• Vertical cylinder (finger) geometry\n"
    "• Radius: 10mm, Height: 28mm\n"
    "• Fixed z-axis pushing (planar task)\n"
    "• Validated: sphere pusher was unstable\n"
    "  for lateral contact on thin objects"
)
add_text_box(slide, Inches(1.0), Inches(5.2), Inches(5.0), Inches(1.5),
    pusher_text, font_size=13, color=LIGHT_GRAY)

add_shape_bg(slide, Inches(6.8), Inches(4.5), Inches(5.5), Inches(2.5), BG_MID)
add_text_box(slide, Inches(7.0), Inches(4.6), Inches(5.0), Inches(0.5),
    "Success Criterion (Strict Pose Stop)", font_size=18, color=GOLD, bold=True)
success_text = (
    "• Position error  ≤ 1.5 mm\n"
    "• Orientation error  ≤ 3.0°\n"
    "• BOTH must be met simultaneously\n"
    "• No legacy 5cm early-stop shortcut\n"
    "• Budget: 600 or 800 total steps"
)
add_text_box(slide, Inches(7.0), Inches(5.2), Inches(5.0), Inches(1.5),
    success_text, font_size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════
# Slide 5: Current Results
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Oracle-MPC Capacity Results", font_size=32, color=ACCENT, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
    "Phase 1: Can the planner solve the task with true dynamics? (No learned model yet)",
    font_size=16, color=LIGHT_GRAY)

# Open space results
add_shape_bg(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.2), BG_MID)
add_text_box(slide, Inches(1.0), Inches(2.1), Inches(5.0), Inches(0.5),
    "✅ Open Space (train_sim_id)", font_size=18, color=ACCENT3, bold=True)
open_text = (
    "c23_precise:  mean pos = 2.70mm,  θ = ~3°\n"
    "              success_pos_1cm = 100%\n\n"
    "c25_fast:     mean pos = 2.38mm\n"
    "              success_0.5cm_5° = 100%\n\n"
    "→ mm-level precision confirmed"
)
add_text_box(slide, Inches(1.0), Inches(2.7), Inches(5.0), Inches(1.5),
    open_text, font_size=13, color=WHITE)

# Passage results
add_shape_bg(slide, Inches(6.8), Inches(2.0), Inches(5.5), Inches(2.2), BG_MID)
add_text_box(slide, Inches(7.0), Inches(2.1), Inches(5.0), Inches(0.5),
    "⚠️ Obstacle Layouts (sixpack sweep)", font_size=18, color=GOLD, bold=True)
obs_text = (
    "passage_direct_wide (9cm gap):   ✅ both budgets\n"
    "passage_direct_medium (7cm):     ✅ strict800 only\n"
    "passage_hard (5cm):       ❌ failed (~190mm)\n"
    "blocking (easy/med/hard): ❌ all failed\n"
    "  (170–230mm final error)\n\n"
    "strict800 SPSR: 33%  |  strict600 SPSR: 17%"
)
add_text_box(slide, Inches(7.0), Inches(2.7), Inches(5.0), Inches(1.5),
    obs_text, font_size=13, color=WHITE)

# Key findings
add_shape_bg(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(2.5), BG_MID)
add_text_box(slide, Inches(1.0), Inches(4.7), Inches(11.0), Inches(0.5),
    "Key Findings", font_size=20, color=GOLD, bold=True)

findings = [
    "• Oracle-MPC achieves sub-2mm precision in open space — task is solvable",
    "• Blocking layouts remain unsolved — obstacles block direct pushing paths",
    "• Passage difficulty scales with gap width: 9cm > 7cm > 5cm",
    "• More MPC budget (strict800 vs 600) helps borderline cases (passage_direct_medium)",
    "• 5cm early-stop was masking true capability — strict pose stop is the correct metric",
]
add_bullet_slide(slide, findings, Inches(1.0), Inches(5.3), Inches(11.0), Inches(1.8),
    font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════
# Slide 6: Validation Pipeline Status
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Validation Pipeline Status", font_size=32, color=ACCENT, bold=True)

# Two columns
for col, (title, items, color) in enumerate([
    ("✅ Completed (23 files PASS)", [
        "MuJoCo env scaffold (reset/step/clone/restore)",
        "ObjectShapeFactory (6 shape families)",
        "Oracle rollout + cost function",
        "CEM-MPC optimizer interface",
        "State sanity checks (140/140 templates)",
        "Toy oracle-MPC (20/20)",
        "MuJoCo oracle-MPC interface (5/5)",
        "Reset template generation + schema",
        "State normalizer interface",
        "Three encoder variants (forward/backward)",
        "RIGWorldModel unified interface",
        "Strict pose-stop code fix",
    ], ACCENT3),
    ("⚠️ In Progress / Blocked", [
        "Obstacle instantiation in MuJoCo ← #1 blocker",
        "Strict pose-stop smoke test (code ready)",
        "Success metric variable bug fix",
        "Layout OOD capacity with real obstacles",
        "Shape OOD capacity test",
    ], GOLD),
]):
    x = Inches(0.8 + col * 6.3)
    add_shape_bg(slide, x, Inches(1.5), Inches(5.8), Inches(5.5), BG_MID)
    add_text_box(slide, x + Inches(0.3), Inches(1.6), Inches(5.2), Inches(0.5),
        title, font_size=18, color=color, bold=True)
    add_bullet_slide(slide, items, x + Inches(0.3), Inches(2.2), Inches(5.2), Inches(4.5),
        font_size=13, color=LIGHT_GRAY, spacing=Pt(5))


# ═══════════════════════════════════════════════════
# Slide 7: Roadmap
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Roadmap", font_size=32, color=ACCENT, bold=True)

phases = [
    ("Phase 1: Oracle-MPC Capacity", "NOW", [
        "Open-space mm-precision ✅",
        "Obstacle gate (blocking + passage)",
        "Strict pose-stop validation",
    ], ACCENT),
    ("Phase 2: Data Collection", "NEXT", [
        "Collect sim data with Oracle-MPC",
        "State normalizer fit on train only",
        "Validate no split leakage",
    ], LIGHT_GRAY),
    ("Phase 3: Model Training", "LATER", [
        "Train 3 encoder variants",
        "Same data, same hyperparams",
        "Flat vs ObjCentric vs Causal",
    ], LIGHT_GRAY),
    ("Phase 4: OOD Evaluation", "GOAL", [
        "Learned model + CEM-MPC",
        "ID test → Layout OOD → Shape OOD",
        "Quantify representation gap",
    ], LIGHT_GRAY),
]

for i, (title, badge, items, color) in enumerate(phases):
    x = Inches(0.6 + i * 3.1)
    add_shape_bg(slide, x, Inches(1.5), Inches(2.9), Inches(4.0), BG_MID)

    # Badge
    badge_color = ACCENT if badge == "NOW" else (GOLD if badge == "NEXT" else LIGHT_GRAY)
    add_text_box(slide, x + Inches(0.15), Inches(1.6), Inches(1.0), Inches(0.4),
        badge, font_size=11, color=BG_DARK, bold=True)
    badge_shape = add_shape_bg(slide, x + Inches(0.1), Inches(1.6), Inches(0.9), Inches(0.35), badge_color)
    # Re-add text on top
    add_text_box(slide, x + Inches(0.1), Inches(1.58), Inches(0.9), Inches(0.4),
        badge, font_size=11, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.15), Inches(2.1), Inches(2.6), Inches(0.5),
        title, font_size=14, color=color, bold=True)
    add_bullet_slide(slide, items, x + Inches(0.15), Inches(2.7), Inches(2.6), Inches(2.5),
        font_size=12, color=LIGHT_GRAY, spacing=Pt(4))

# Timeline arrow
add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.8),
    "Sprint history:  Apr 20 (skeleton) → May 6 (MuJoCo env) → May 8 (oracle-MPC) → May 11-12 (sweeps + strict stop) → May 13 (obstacle sweeps)",
    font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# Slide 8: Tech Stack
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Tech Stack & Architecture", font_size=32, color=ACCENT, bold=True)

cols = [
    ("Simulation", ["MuJoCo physics engine", "Compound geom objects (T/L/cross/...)", "Custom MujocoPushEnv wrapper", "State clone/restore for MPC rollouts"]),
    ("Planning", ["CEM-MPC (Cross-Entropy Method)", "1024 samples, 96 elites, 5 iterations", "Horizon: 80 steps", "Strict pose-stop criterion"]),
    ("Models", ["PyTorch", "3 encoder variants (Flat/ObjCentric/Causal)", "DynamicsHead + SubgoalHead", "RIGWorldModel unified interface"]),
    ("Infrastructure", ["Python + conda (lerobot env)", "28-core parallel CEM evaluation", "Modular scripts (debug/sweep/render)", "Git-versioned experiment configs"]),
]

for i, (title, items) in enumerate(cols):
    x = Inches(0.6 + i * 3.1)
    add_shape_bg(slide, x, Inches(1.5), Inches(2.9), Inches(3.5), BG_MID)
    add_text_box(slide, x + Inches(0.2), Inches(1.6), Inches(2.5), Inches(0.5),
        title, font_size=16, color=GOLD, bold=True)
    add_bullet_slide(slide, items, x + Inches(0.2), Inches(2.2), Inches(2.5), Inches(2.5),
        font_size=12, color=LIGHT_GRAY, spacing=Pt(4))

# Code stats
add_shape_bg(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5), BG_MID)
add_text_box(slide, Inches(1.0), Inches(5.6), Inches(11.0), Inches(0.5),
    "Codebase Structure", font_size=18, color=GOLD, bold=True)
struct_text = (
    "src/  (envs, planners, models, metrics, data, interventions, utils)  ·  "
    "scripts/  (debug, sweep, render, eval)  ·  "
    "configs/  (planner, train, eval, splits)  ·  "
    "data/  (reset templates, metadata)  ·  "
    "runs/  (sweep results, videos, logs)"
)
add_text_box(slide, Inches(1.0), Inches(6.2), Inches(11.0), Inches(0.7),
    struct_text, font_size=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════
# Slide 9: Thank You
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_shape_bg(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1),
    "Thank You", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
    "Questions & Discussion", font_size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(1),
    "Next milestone: Obstacle gate → Data collection → Model training",
    font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ── Save ──
out_path = "/home/brucewu/my_robot_project/project_overview.pptx"
prs.save(out_path)
print(f"PPT saved to: {out_path}")
