#!/usr/bin/env python3
"""
组会分享 PPT：因果感知世界模型与机器人操控
—— 从 C-JEPA 到物体级推放任务的 OOD 泛化

10 页，深色主题学术风格
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 颜色方案 (深色学术风) ──
BG_DARK    = RGBColor(0x0D, 0x11, 0x17)   # 深黑蓝背景
BG_ACCENT  = RGBColor(0x14, 0x1B, 0x25)   # 稍亮卡片
TEXT_WHITE = RGBColor(0xE8, 0xEA, 0xED)   # 正文白
TEXT_GRAY  = RGBColor(0x9A, 0xA0, 0xA6)   # 次级灰
ACCENT_BLUE  = RGBColor(0x4D, 0xAB, 0xF7) # 亮蓝强调
ACCENT_GREEN = RGBColor(0x34, 0xD3, 0x99) # 绿色
ACCENT_ORANGE= RGBColor(0xFF, 0xA7, 0x26) # 橙色
ACCENT_RED   = RGBColor(0xFF, 0x52, 0x52) # 红色
ACCENT_PURPLE= RGBColor(0xBB, 0x86, 0xFC) # 紫色
BORDER_SUBTLE = RGBColor(0x30, 0x3A, 0x48) # 边框

prs = Presentation()
prs.slide_width  = Inches(13.333)  # 16:9 宽屏
prs.slide_height = Inches(7.5)

# ── 辅助函数 ──
def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_card(slide, left, top, width, height, fill=BG_ACCENT, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=TEXT_WHITE, spacing=Pt(8), font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0 if not item.startswith("  ") else 1
        # bullet symbol
        if item.startswith("  "):
            p.text = "    ▸ " + item.strip()
        else:
            p.text = "● " + item
    return txBox

def add_section_title(slide, text, y=Inches(0.3)):
    add_text_box(slide, Inches(0.8), y, Inches(11), Inches(0.6), text,
                 font_size=28, color=ACCENT_BLUE, bold=True)

def add_page_number(slide, num, total=10):
    add_text_box(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4),
                 f"{num}/{total}", font_size=11, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)

def make_flow_arrow(slide, left, top, width):
    """Draw a simple right-arrow"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.25))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    return shape


# ══════════════════════════════════════════════════════════════
# PAGE 1: 标题页
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# 装饰线
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.8), Inches(1.5), Pt(4))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT_BLUE; line.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(1.2),
             "因果感知世界模型与机器人操控",
             font_size=40, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(0.8),
             "从 C-JEPA 到物体级推放任务的 OOD 泛化",
             font_size=26, color=ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.6),
             "Causal World Models meets Robot Manipulation: Object-Level OOD Generalization in Pushing Tasks",
             font_size=14, color=TEXT_GRAY)

add_text_box(slide, Inches(0.8), Inches(5.5), Inches(6), Inches(0.5),
             "组会分享 ｜ 2026.05",
             font_size=16, color=TEXT_GRAY)

# 两个项目标签
add_card(slide, Inches(0.8), Inches(6.2), Inches(3.5), Inches(0.7), fill=BG_ACCENT, border_color=ACCENT_BLUE)
add_text_box(slide, Inches(1.0), Inches(6.3), Inches(3.1), Inches(0.5),
             "🏷 C-JEPA · 因果世界模型", font_size=13, color=ACCENT_BLUE)
add_card(slide, Inches(4.7), Inches(6.2), Inches(4.0), Inches(0.7), fill=BG_ACCENT, border_color=ACCENT_GREEN)
add_text_box(slide, Inches(4.9), Inches(6.3), Inches(3.6), Inches(0.5),
             "🏷 Robot Pushing · 推放 OOD 泛化", font_size=13, color=ACCENT_GREEN)

add_page_number(slide, 1)

# ══════════════════════════════════════════════════════════════
# PAGE 2: Motivation — 为什么需要真正的世界模型？
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "Motivation：为什么需要真正的世界模型？")

# 问题对比: 三列
col_w = Inches(3.6)
gap = Inches(0.35)
start_x = Inches(0.8)

titles = ["LLM 霸权幻觉", "VLA 的根本局限", "OpenClaw 的启示"]
bodies = [
    "● LLM 太成功 → 领域误以为\n   \"语言理解\" = \"物理理解\"\n\n● 文本统计 ≠ 因果结构\n   • 杯子会碎 ≠ 知道怎么推\n   • 水会洒 ≠ 知道洒多少\n\n● 资金、人才涌向 LLM-based\n   路线，世界模型被边缘化",
    "● VLA = 视觉 token + 动作 token\n   → 全部压扁成 seq2seq\n\n● 核心矛盾：\n   行动改变世界状态，\n   语言不改变。\n\n● 行动有：物理后果、不可逆性、\n   safety 约束 — 这些不是\n   token prediction 能解决的",
    "● OpenClaw 证明：\n   预训练大脑 + 工具接口\n   = 功能完备的 agent\n\n● 推广到物理世界：\n   预训练世界模型\n   + 机器人身体\n   = 通用机器人\n\n● 缺的不是身体，是大脑。"
]

for i, (t, b) in enumerate(zip(titles, bodies)):
    x = start_x + i * (col_w + gap)
    add_card(slide, x, Inches(1.3), col_w, Inches(5.5), fill=BG_ACCENT, border_color=[ACCENT_ORANGE, ACCENT_RED, ACCENT_GREEN][i])
    add_text_box(slide, x + Inches(0.25), Inches(1.45), col_w - Inches(0.5), Inches(0.5),
                 t, font_size=17, color=[ACCENT_ORANGE, ACCENT_RED, ACCENT_GREEN][i], bold=True)
    add_text_box(slide, x + Inches(0.25), Inches(2.0), col_w - Inches(0.5), Inches(4.5),
                 b, font_size=12, color=TEXT_WHITE)

# 核心观点
add_card(slide, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.35), fill=BG_ACCENT, border_color=ACCENT_BLUE)
add_text_box(slide, Inches(1.0), Inches(7.02), Inches(11.3), Inches(0.3),
             "核心立场：\"真正通向通用机器人的路，不在 LLM 的延长线上，而在世界模型的另一侧。\" —— 谢赛宁",
             font_size=12, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 2)

# ══════════════════════════════════════════════════════════════
# PAGE 3: C-JEPA 核心思想
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "C-JEPA：通过对象级 Masking 诱导因果归纳偏置")

# 核心公式
add_card(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.6), fill=BG_ACCENT, border_color=ACCENT_PURPLE)
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(5.1), Inches(0.4),
             "核心创新", font_size=18, color=ACCENT_PURPLE, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(1.85), Inches(5.1), Inches(2.0), [
    "从 image-patch 级 JEPA → 对象级 JEPA",
    "训练时 mask 部分 object slot",
    "  被 mask 的 slot 必须从其他 slot 推断",
    "   → 强制学习对象间交互关系",
    "   → 产生 Counterfactual-like 效果",
    "   → 防止捷径解 (shortcut solutions)",
    "关键等价：Object-Level Masking",
    "  = Latent Intervention（潜在干预）",
], font_size=12, color=TEXT_WHITE)

add_card(slide, Inches(6.7), Inches(1.3), Inches(5.8), Inches(2.6), fill=BG_ACCENT, border_color=ACCENT_PURPLE)
add_text_box(slide, Inches(6.9), Inches(1.4), Inches(5.4), Inches(0.4),
             "形式化分析", font_size=18, color=ACCENT_PURPLE, bold=True)
add_bullet_list(slide, Inches(6.9), Inches(1.85), Inches(5.4), Inches(2.0), [
    "传统 JEPA: x → f(x) → predictor → f(x+)",
    "  预测未来，但不要求因果推理",
    "",
    "C-JEPA: x → f(x) → [mask object k]",
    "        → predictor(args from other obj.)",
    "        → f_k(x+)（被 mask 对象的状态）",
    "",
    "Masking = 切断某个对象的因果链",
    "→ predictor 必须从其他对象推断",
    "   → 隐式地学习 P(obj_k | do(other_objs))",
], font_size=12, color=TEXT_WHITE)

# 三个环境
add_card(slide, Inches(0.8), Inches(4.1), Inches(11.7), Inches(3.0), fill=BG_ACCENT, border_color=BORDER_SUBTLE)
add_text_box(slide, Inches(1.0), Inches(4.2), Inches(11.3), Inches(0.4),
             "实验环境 & 对象级表征", font_size=18, color=ACCENT_BLUE, bold=True)

envs = [
    ("CLEVRER", ACCENT_ORANGE,
     "● 碰撞事件视觉推理\n● VQA: descriptive / explanatory / predictive / counterfactual\n● VideoSAUR / SAVi encoder → 7 slots × 128d"),
    ("PushT", ACCENT_GREEN,
     "● 块推送控制任务\n● 蓝圈 agent + 灰T目标 + 绿T干扰物\n● VideoSAUR encoder → 4 slots × 128d\n● CEM-MPC latent space planning"),
    ("PHYRE", ACCENT_BLUE,
     "● 物理推理游戏\n● 多对象互动场景\n● 验证对象级理解泛化"),
]

for i, (name, color, desc) in enumerate(envs):
    x = Inches(1.0) + i * Inches(3.9)
    add_card(slide, x, Inches(4.7), Inches(3.6), Inches(2.2), fill=BG_DARK, border_color=color)
    add_text_box(slide, x + Inches(0.15), Inches(4.8), Inches(3.3), Inches(0.35),
                 name, font_size=15, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), Inches(5.2), Inches(3.3), Inches(1.5),
                 desc, font_size=11, color=TEXT_WHITE)

add_page_number(slide, 3)

# ══════════════════════════════════════════════════════════════
# PAGE 4: C-JEPA 架构详解
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "C-JEPA 架构详解")

# 架构流程: 水平 flow
flow_items = [
    ("观测图像\n(224×224)", ACCENT_BLUE),
    ("VideoSAUR\nSlot Attention", ACCENT_PURPLE),
    ("4 Slots\n(4×128 dim)", ACCENT_ORANGE),
    ("Object-Level\nMasking", ACCENT_RED),
    ("Non-Causal\nTransformer", ACCENT_GREEN),
    ("Predicted\nFuture Slots", ACCENT_BLUE),
]

x = Inches(0.5)
for i, (label, color) in enumerate(flow_items):
    w = Inches(1.8)
    add_card(slide, x, Inches(1.3), w, Inches(1.1), fill=BG_ACCENT, border_color=color)
    add_text_box(slide, x + Inches(0.1), Inches(1.4), w - Inches(0.2), Inches(0.9),
                 label, font_size=11, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    if i < len(flow_items) - 1:
        make_flow_arrow(slide, x + w + Inches(0.05), Inches(1.7), Inches(0.2))
    x += w + Inches(0.3)

# 三个核心机制
mechanisms = [
    ("Identity Anchor", ACCENT_PURPLE,
     "● t=0 的 slot 通过 id_projector\n  投影为 query vector\n● 作为对象身份的\"锚点\"\n● \"这个对象一开始长这样，\n  预测它以后会怎样\"\n● 关键技术: 时间位置编码\n  + anchor query + mask token"),
    ("Hungarian Matching", ACCENT_ORANGE,
     "● Slot 是无序集合 → 直接 MSE\n  因排列不同会得到错误高 cost\n● 匈牙利算法: 构建 cost[i][j],\n  找最优 slot-to-slot 对应\n● 训练: hungarian_matching_loss\n● 推演: reorder_slots_to_match\n● 规划: hungarian_cost (评价用)\n● 保证 slot ID 在时间上一致"),
    ("Non-Causal Transformer", ACCENT_GREEN,
     "● 所有 token 互相可 attend\n  (bidirectional, non-autoregressive)\n● 训练: 部分 slot 被 mask,\n  从未 mask slot 推断被 mask 的\n● 推理: 所有 slot 可见,\n  未来位置放 query token\n● 核心: Transformer 从 history\n  提取信息预测未来"),
]

for i, (title, color, desc) in enumerate(mechanisms):
    x = Inches(0.8) + i * Inches(4.2)
    add_card(slide, x, Inches(2.7), Inches(3.9), Inches(4.2), fill=BG_ACCENT, border_color=color)
    add_text_box(slide, x + Inches(0.15), Inches(2.8), Inches(3.6), Inches(0.35),
                 title, font_size=15, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), Inches(3.2), Inches(3.6), Inches(3.5),
                 desc, font_size=11, color=TEXT_WHITE)

add_page_number(slide, 4)

# ══════════════════════════════════════════════════════════════
# PAGE 5: C-JEPA 关键结果
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "C-JEPA 关键实验结果")

# 结果卡片
add_card(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.4), fill=BG_ACCENT, border_color=ACCENT_GREEN)
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(5.1), Inches(0.4),
             "🏆 CLEVRER 视觉问答", font_size=18, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(1.9), Inches(5.1), Inches(1.6), [
    "反事实推理 (counterfactual): +20% 绝对提升",
    "  同一架构，仅加 object-level masking",
    "其余推理类型 (descriptive/explanatory/",
    "  predictive) 保持或提升",
    "对象级 masking 选择性地提升因果推理能力",
    '  而非简单的"更多数据 = 更好"',
], font_size=12, color=TEXT_WHITE)

add_card(slide, Inches(6.7), Inches(1.3), Inches(5.8), Inches(2.4), fill=BG_ACCENT, border_color=ACCENT_BLUE)
add_text_box(slide, Inches(6.9), Inches(1.4), Inches(5.4), Inches(0.4),
             "🏆 PushT 操控规划", font_size=18, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(6.9), Inches(1.9), Inches(5.4), Inches(1.6), [
    "Latent 维度仅为 patch-based 的 ~1%",
    "  4 slots × 128d = 512d vs 196×384d = 75Kd",
    "CEM-MPC: 300×30=9000 次前向时差距显著",
    "  规划速度大幅提升，性能可比",
    "无需重建像素: 直接在 latent space 中",
    "  推演、评价、优化动作",
], font_size=12, color=TEXT_WHITE)

# 对比表
add_card(slide, Inches(0.8), Inches(3.9), Inches(11.7), Inches(1.8), fill=BG_ACCENT, border_color=BORDER_SUBTLE)
add_text_box(slide, Inches(1.0), Inches(4.0), Inches(5), Inches(0.4),
             "Latent Space 效率对比", font_size=16, color=TEXT_WHITE, bold=True)

table_data = [
    ["方法", "Latent 维度", "每步计算", "规划效率", "性能"],
    ["Patch-based (Dreamer)", "196×384 = ~75K", "高", "慢", "Baseline"],
    ["C-JEPA (ours)", "4×128 = 512", "极低", "快", "可比 / 更优"],
]
for row_i, row in enumerate(table_data):
    for col_i, cell in enumerate(row):
        c = ACCENT_BLUE if row_i == 0 else (ACCENT_GREEN if row_i == 2 else TEXT_WHITE)
        b = True if row_i == 0 else False
        add_text_box(slide, Inches(1.0) + col_i * Inches(2.8), Inches(4.45) + row_i * Inches(0.4),
                     Inches(2.6), Inches(0.35), cell, font_size=12, color=c, bold=b)

add_card(slide, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.2), fill=BG_ACCENT, border_color=ACCENT_ORANGE)
add_text_box(slide, Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.4),
             "💡 关键洞察", font_size=16, color=ACCENT_ORANGE, bold=True)
add_text_box(slide, Inches(1.0), Inches(6.35), Inches(11.3), Inches(0.6),
             "对象级表征 × 因果 masking = 既高效（~1% latent dim）又强泛化（+20% counterfactual）的世界模型。这为机器人操控提供了轻量级但强大的规划基础。",
             font_size=12, color=TEXT_WHITE)

add_page_number(slide, 5)

# ══════════════════════════════════════════════════════════════
# PAGE 6: 机器人推放平台 & Benchmark
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "机器人推放实验平台 & Benchmark 设计")

# 左边: 实验设置
add_card(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.7), fill=BG_ACCENT, border_color=ACCENT_BLUE)
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(5.1), Inches(0.4),
             "实验设置", font_size=18, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(1.9), Inches(5.1), Inches(1.8), [
    "仿真环境: MuJoCo + SO-101 机械臂",
    "任务: 结构化对象推放 (Object Pushing)",
    "对象: L 形、T 形刚体",
    "动作为 2D 推杆末端位移 (x, y)",
    "固定 CEM-MPC 规划器",
    "物体级结构状态 (object states)",
    "  非 RGB 像素 → 因果感知高层表征",
], font_size=12, color=TEXT_WHITE)

add_text_box(slide, Inches(1.0), Inches(3.6), Inches(5.1), Inches(0.4),
             "两阶段研究路径", font_size=16, color=ACCENT_PURPLE, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(4.0), Inches(5.1), Inches(2.5), [
    "Stage 1 (当前): Oracle → 验证 MPC 能力上限",
    "  → MuJoCo oracle rollout",
    "  → Oracle-MPC capacity 评估",
    "Stage 2: Learned Model → 测量 OOD 距离",
    "  → sim 数据收集",
    "  → 训练 high-level 因果感知模型",
    "  → Learned model + MPC → OOD 泛化 gap",
], font_size=12, color=TEXT_WHITE)

# 右边: Benchmark
add_card(slide, Inches(6.7), Inches(1.3), Inches(5.8), Inches(5.7), fill=BG_ACCENT, border_color=ACCENT_GREEN)
add_text_box(slide, Inches(6.9), Inches(1.4), Inches(5.4), Inches(0.4),
             "Benchmark 体系", font_size=18, color=ACCENT_GREEN, bold=True)

bench_items = [
    ("🏗 任务 (tasks.md)", "Narrow Passage / Edge Goal /\nBlocking / Open Space / Mixed"),
    ("📊 指标 (metrics.md)", "Success Rate / Completion Rate /\nSteps to Goal / Path Efficiency"),
    ("📁 数据 Splits", "train_sim_id / test_sim_id (layout OOD)\n/ test_shape_id (shape OOD)"),
    ("🏷 模板系统", "blocking / edge_goal /\nnarrow_passage / open_space"),
    ("📋 基准 (baselines.md)", "Oracle MPC / Random / Heuristic /\nLearned Model"),
]

y = Inches(1.9)
for title, desc in bench_items:
    add_text_box(slide, Inches(6.9), y, Inches(5.4), Inches(0.3),
                 title, font_size=13, color=ACCENT_GREEN, bold=True)
    add_text_box(slide, Inches(7.3), y + Inches(0.3), Inches(5.0), Inches(0.7),
                 desc, font_size=11, color=TEXT_GRAY)
    y += Inches(0.95)

add_page_number(slide, 6)

# ══════════════════════════════════════════════════════════════
# PAGE 7: OOD 泛化设计
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "OOD 泛化：Layout & Shape 两个维度")

# Layout OOD
add_card(slide, Inches(0.8), Inches(1.3), Inches(5.7), Inches(2.8), fill=BG_ACCENT, border_color=ACCENT_ORANGE)
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(5.3), Inches(0.4),
             "Layout OOD（主维度）", font_size=18, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, Inches(1.0), Inches(1.9), Inches(5.3), Inches(2.0), [
    "训练分布 (train_sim_id):",
    "  特定物体排列的训练模板",
    "测试分布 (test_sim_id):",
    "  未见过的物体布局配置",
    "挑战: 模型必须理解物体间空间关系",
    "  \"T 形在 L 形左侧 10cm\" →",
    "  \"T 形在 L 形右侧 15cm\" 的泛化",
    "本质: 测试模型是否学会了",
    "  结构化的物体关系而非记忆位置",
], font_size=12, color=TEXT_WHITE)

# Shape OOD
add_card(slide, Inches(6.9), Inches(1.3), Inches(5.6), Inches(2.8), fill=BG_ACCENT, border_color=ACCENT_GREEN)
add_text_box(slide, Inches(7.1), Inches(1.4), Inches(5.2), Inches(0.4),
             "Shape OOD（次维度）", font_size=18, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(7.1), Inches(1.9), Inches(5.2), Inches(2.0), [
    "训练形状 vs 测试形状",
    "  如: 训练见 L/T 形，测试新形态",
    "挑战: 模型必须从形状抽象出",
    "  物理属性（质心、惯性、接触面）",
    "更接近真实世界的泛化需求:",
    "  机器人遇到新物体时不能重训练",
    "因果表征的优势:",
    "  物体属性是因果变量 →",
    "  形状变化不改变因果结构",
], font_size=12, color=TEXT_WHITE)

# 模板系统
add_card(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.7), fill=BG_ACCENT, border_color=BORDER_SUBTLE)
add_text_box(slide, Inches(1.0), Inches(4.4), Inches(11.3), Inches(0.4),
             "模板系统 (Template System) — 结构化 OOD 生成", font_size=16, color=ACCENT_PURPLE, bold=True)

templates = [
    ("🔄 Blocking", "阻挡", "物体阻挡\n目标路径"),
    ("⏺ Edge Goal", "边角目标", "目标靠近\n环境边界"),
    ("⏳ Narrow Passage", "窄通道", "物体必须\n穿过窄缝"),
    ("⬜ Open Space", "开放空间", "无障碍物\n自由推放"),
    ("🔀 Mixed", "混合", "多元素\n组合场景"),
]

for i, (emoji_title, _, desc) in enumerate(templates):
    x = Inches(1.2) + i * Inches(2.3)
    add_card(slide, x, Inches(4.9), Inches(2.0), Inches(1.9), fill=BG_DARK, border_color=BORDER_SUBTLE)
    add_text_box(slide, x + Inches(0.1), Inches(5.0), Inches(1.8), Inches(0.35),
                 emoji_title, font_size=13, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(5.4), Inches(1.8), Inches(1.2),
                 desc, font_size=11, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 7)

# ══════════════════════════════════════════════════════════════
# PAGE 8: 当前进展
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "当前实验进展")

# 三列进展
progress_items = [
    ("✅ 已完成", ACCENT_GREEN, [
        "MuJoCo 仿真环境搭建 (SO-101)",
        "L/T 形物体模型构建",
        "CEM-MPC 规划器集成",
        "Oracle rollout & MPC capacity 评估",
        "5 类模板生成系统",
        "结构化 benchmark (指标/分片/基准)",
        "Layout OOD 的初步 MPC 成功率",
        "Shape OOD 的初步分析",
    ]),
    ("🔄 进行中", ACCENT_ORANGE, [
        "Oracle-MPC 在 OOD split 上的\n  完整评估 (boundary analysis)",
        "各模板类型的困难度定性",
        "数据收集 pipeline (sim → dataset)",
        "High-level 因果感知模型\n  架构设计",
    ]),
    ("📋 待开展", ACCENT_BLUE, [
        "训练 causality-aware high-level model",
        "Learned model + MPC → OOD gap 量化",
        "与 C-JEPA 的衔接:",
        "  对象 slot → 世界模型 → 规划",
        "多 template 上的全面对比",
        "论文撰写",
    ]),
]

for i, (title, color, items) in enumerate(progress_items):
    x = Inches(0.8) + i * Inches(4.2)
    add_card(slide, x, Inches(1.3), Inches(3.9), Inches(5.8), fill=BG_ACCENT, border_color=color)
    add_text_box(slide, x + Inches(0.15), Inches(1.4), Inches(3.6), Inches(0.4),
                 title, font_size=17, color=color, bold=True)
    item_text = "\n".join(f"  ▸ {it}" for it in items)
    add_text_box(slide, x + Inches(0.15), Inches(1.9), Inches(3.6), Inches(5.0),
                 item_text, font_size=11, color=TEXT_WHITE)

add_page_number(slide, 8)

# ══════════════════════════════════════════════════════════════
# PAGE 9: 端到端管线 & 未来路径
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_section_title(slide, "端到端管线：从 C-JEPA 到真实机器人操控")

# 主流程图
pipeline_items = [
    ("观测图像", "RGB camera\n→ raw pixels", ACCENT_BLUE),
    ("对象 Slot\n编码器", "VideoSAUR / SAVi\n→ object slots", ACCENT_PURPLE),
    ("C-JEPA\n世界模型", "学习因果转移\nslot_t → slot_{t+1}", ACCENT_ORANGE),
    ("CEM-MPC\n规划器", "Latent space\nplanning & control", ACCENT_GREEN),
    ("机器人\n执行", "SO-101 / real\nrobot action", ACCENT_RED),
]

x = Inches(0.5)
for i, (label, desc, color) in enumerate(pipeline_items):
    w = Inches(2.1)
    add_card(slide, x, Inches(1.4), w, Inches(1.6), fill=BG_ACCENT, border_color=color)
    add_text_box(slide, x + Inches(0.1), Inches(1.5), w - Inches(0.2), Inches(0.5),
                 label, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(2.05), w - Inches(0.2), Inches(0.8),
                 desc, font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
    if i < len(pipeline_items) - 1:
        make_flow_arrow(slide, x + w + Inches(0.05), Inches(2.05), Inches(0.2))
    x += w + Inches(0.3)

# 三个关键问题
add_card(slide, Inches(0.8), Inches(3.3), Inches(11.7), Inches(3.5), fill=BG_ACCENT, border_color=BORDER_SUBTLE)
add_text_box(slide, Inches(1.0), Inches(3.4), Inches(11.3), Inches(0.4),
             "关键研究问题 & 未来方向", font_size=18, color=ACCENT_BLUE, bold=True)

questions = [
    ("Q1: 表征迁移", ACCENT_ORANGE,
     "C-JEPA 在 CLEVRER/PushT 上学到的对象级表示\n能否迁移到 MuJoCo SO-101 推放任务？\n→ 关键挑战: sim-to-sim domain gap, object appearance 差异"),
    ("Q2: OOD 泛化边界", ACCENT_GREEN,
     "因果世界模型在 OOD layout/shape 上的泛化上限在哪？\n→ 量化 causal inductive bias 的实际收益\n→ 与纯数据驱动的 Imitation/RL baseline 对比"),
    ("Q3: 真实机器人部署", ACCENT_PURPLE,
     "从 MuJoCo → 真实 SO-101 的 sim-to-real gap\n→ 对象 slot 的鲁棒性: 能否适应真实感知噪声？\n→ 闭环 replanning 的实时性"),
]

y = Inches(3.85)
for title, color, desc in questions:
    add_text_box(slide, Inches(1.0), y, Inches(11.3), Inches(0.3),
                 title, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(1.2), y + Inches(0.32), Inches(11.1), Inches(0.75),
                 desc, font_size=11, color=TEXT_WHITE)
    y += Inches(1.05)

add_page_number(slide, 9)

# ══════════════════════════════════════════════════════════════
# PAGE 10: 总结
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# 装饰线
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(1.5), Pt(4))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT_BLUE; line.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "总结与展望", font_size=36, color=TEXT_WHITE, bold=True)

# 三个 takeaway
takeaways = [
    ("🎯 因果世界模型 > 统计世界模型",
     ACCENT_BLUE,
     "C-JEPA 通过 object-level masking 在 latent space 中诱导因果归纳偏置。\n不依赖 LLM，不依赖像素重建，以 ~1% latent dimension 实现可比规划性能。"),
    ("🤖 OOD 泛化需要结构化理解",
     ACCENT_GREEN,
     "推放任务中的 layout/shape OOD 不是 scale data 能解决的。\n需要模型理解\"物体是什么\"、\"在哪\"、\"如何互动\"——这是因果表征的核心。"),
    ("🧠 大脑 > 身体",
     ACCENT_ORANGE,
     "OpenClaw 的启示：预训练大脑 + 身体接口 = 通用 agent。\n机器人领域不应只卷身体和数据，更需要真正理解物理因果的世界模型。"),
]

y = Inches(1.6)
for title, color, desc in takeaways:
    add_card(slide, Inches(0.8), y, Inches(11.7), Inches(1.45), fill=BG_ACCENT, border_color=color)
    add_text_box(slide, Inches(1.0), y + Inches(0.1), Inches(11.3), Inches(0.4),
                 title, font_size=17, color=color, bold=True)
    add_text_box(slide, Inches(1.0), y + Inches(0.5), Inches(11.3), Inches(0.8),
                 desc, font_size=12, color=TEXT_WHITE)
    y += Inches(1.6)

# 底部项目链接
add_card(slide, Inches(0.8), Inches(6.5), Inches(5.5), Inches(0.6), fill=BG_ACCENT, border_color=BORDER_SUBTLE)
add_text_box(slide, Inches(1.0), Inches(6.6), Inches(5.1), Inches(0.4),
             "📎 C-JEPA: arxiv.org/abs/2602.11389  |  hazel-heejeong-nam.github.io/cjepa",
             font_size=11, color=TEXT_GRAY)
add_card(slide, Inches(6.7), Inches(6.5), Inches(5.8), Inches(0.6), fill=BG_ACCENT, border_color=BORDER_SUBTLE)
add_text_box(slide, Inches(6.9), Inches(6.6), Inches(5.4), Inches(0.4),
             "📎 Robot Pushing: /home/brucewu/my_robot_project",
             font_size=11, color=TEXT_GRAY)

add_page_number(slide, 10)

# ── 保存 ──
output_dir = "/home/brucewu/my_robot_project/artifacts"
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "group_meeting_20260514.pptx")
prs.save(output_path)
print(f"✅ PPT saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
